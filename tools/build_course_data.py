from __future__ import annotations

import json
from pathlib import Path

import pdfplumber
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]


LECTURE_SOURCE = Path(r"E:\College\term_6\System modelling\Simulation.pdf")
SECTION_SOURCES = [
    {
        "id": "section-1",
        "title": "Section 1: Simulation Foundations",
        "kind": "pptx",
        "source": Path(r"E:\College\term_6\System modelling\section 1\sim week 1 sec .pptx"),
        "material": "materials/sim-week-1-section.pptx",
        "theme": "Introductory section deck covering simulation meaning, grading, and why simulation is useful.",
    },
    {
        "id": "section-2",
        "title": "Section 2: Simulation Exercises",
        "kind": "pdf",
        "source": Path(r"E:\College\term_6\System modelling\section 2\Simulation week 2 sec.pdf"),
        "material": "materials/simulation-week-2-section.pdf",
        "theme": "Scanned section sheet for early simulation exercises.",
    },
    {
        "id": "section-3",
        "title": "Section 3: Newsboy Problem",
        "kind": "pdf",
        "source": Path(r"E:\College\term_6\System modelling\section 3\Simulation-Section_3.pdf"),
        "material": "materials/simulation-section-3.pdf",
        "theme": "Demand distributions, random-number mapping, and profit analysis for the newsboy model.",
    },
    {
        "id": "accident-problem",
        "title": "Section 4: Accident Penalty Problem",
        "kind": "pdf",
        "source": Path(r"E:\College\term_6\System modelling\section 4\Accident Problem.pdf"),
        "material": "materials/accident-problem.pdf",
        "theme": "Daily accident occurrence and penalty simulation over 12 days.",
    },
    {
        "id": "football-problem",
        "title": "Section 4: Football Injury Problem",
        "kind": "pdf",
        "source": Path(r"E:\College\term_6\System modelling\section 4\Football problem.pdf"),
        "material": "materials/football-problem.pdf",
        "theme": "Random major/minor injuries and running-back stock over a 10-game season.",
    },
    {
        "id": "section-5",
        "title": "Section 5: Mean, Variance, and CDF",
        "kind": "pdf",
        "source": Path(r"E:\College\term_6\System modelling\section 5\Simulation-Section_5.pdf"),
        "material": "materials/simulation-section-5.pdf",
        "theme": "Expected value, variance, discrete distributions, CDF derivation, and random observations.",
    },
    {
        "id": "section-5-scanned",
        "title": "Section 5: Additional Scanned Sheet",
        "kind": "pdf",
        "source": Path(r"E:\College\term_6\System modelling\section 5\Section (5). simulation.pdf"),
        "material": "materials/section-5-simulation-scanned.pdf",
        "theme": "Original scanned section material preserved as the source file.",
    },
    {
        "id": "section-6-scanned",
        "title": "Section 6: Scanned Simulation Sheet",
        "kind": "pdf",
        "source": Path(r"E:\College\term_6\System modelling\section 6\Section (6). Simulation.pdf"),
        "material": "materials/section-6-simulation-scanned.pdf",
        "theme": "Original scanned section material preserved as the source file.",
    },
    {
        "id": "section-7",
        "title": "Section 7: Distribution-Based Random Observations",
        "kind": "pdf",
        "source": Path(r"E:\College\term_6\System modelling\section 7\Simulation-Section_7.pdf"),
        "material": "materials/simulation-section-7.pdf",
        "theme": "Uniform and exponential CDF derivations plus mixed-distribution random observations.",
    },
]

MODULES = [
    {
        "id": "concepts",
        "title": "1. Basic Concepts and Terminology",
        "pages": [6, 26],
        "topics": [
            "System definition and state",
            "System classification",
            "System methodology",
            "Validation and verification",
            "Advantages and disadvantages of simulation",
            "Introductory simulation examples",
        ],
    },
    {
        "id": "probability",
        "title": "2. Probability and Distribution Theory",
        "pages": [27, 43],
        "topics": [
            "Random variables",
            "Events and set operations",
            "Discrete and continuous distributions",
            "CDF and PDF properties",
            "Expected value, variance, and moments",
            "Bernoulli, uniform, Poisson, exponential, and normal distributions",
        ],
    },
    {
        "id": "variance",
        "title": "3. Estimation and Statistical Ideas",
        "pages": [44, 47],
        "topics": [
            "Variance reduction",
            "Complement random numbers",
            "Empirical estimation examples",
            "Random-variable transformations",
        ],
    },
    {
        "id": "queueing",
        "title": "4. Introduction to Queuing Theory",
        "pages": [48, 53],
        "topics": [
            "Customer arrivals and service",
            "Queue length and waiting time",
            "M/M/1/∞/FIFO system",
            "Little's formula",
            "Performance measures",
        ],
    },
    {
        "id": "random-generation",
        "title": "5. Random Numbers and Random Observations",
        "pages": [54, 61],
        "topics": [
            "Pseudo-random numbers",
            "Congruential generators",
            "Allocation method",
            "Inverse transform method",
            "Insurance sales simulation",
        ],
    },
    {
        "id": "discrete-simulation",
        "title": "6. Discrete-System Simulation",
        "pages": [62, 73],
        "topics": [
            "Simulation time and clock management",
            "Time-scan method",
            "Event-scan method",
            "Manufacturing and M/M/1 examples",
            "Collecting and recording simulation data",
        ],
    },
]


def clean_text(text: str) -> str:
    return "\n".join(line.strip() for line in text.replace("\r", "").splitlines() if line.strip())


def extract_pdf_pages(path: Path) -> list[dict]:
    pages = []
    with pdfplumber.open(path) as pdf:
        for index, page in enumerate(pdf.pages, 1):
            text = clean_text(page.extract_text() or "")
            pages.append({"page": index, "text": text, "wordCount": len(text.split())})
    return pages


def extract_pptx_slides(path: Path) -> list[dict]:
    presentation = Presentation(path)
    slides = []
    for index, slide in enumerate(presentation.slides, 1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(clean_text(shape.text))
        text = clean_text("\n".join(parts))
        slides.append({"slide": index, "text": text, "wordCount": len(text.split())})
    return slides


def sample_text(items: list[dict], key: str = "text", limit: int = 420) -> str:
    text = " ".join(item[key] for item in items if item.get(key))
    return text[:limit].strip()


def main() -> None:
    lecture_pages = extract_pdf_pages(LECTURE_SOURCE)
    for module in MODULES:
        start, end = module["pages"]
        module_pages = [page for page in lecture_pages if start <= page["page"] <= end]
        module["pageRefs"] = [{"page": p["page"], "text": p["text"]} for p in module_pages]
        module["summary"] = sample_text(module_pages)

    sections = []
    for item in SECTION_SOURCES:
        if item["kind"] == "pptx":
            entries = extract_pptx_slides(item["source"])
            content_label = "slides"
        else:
            entries = extract_pdf_pages(item["source"])
            content_label = "pages"

        searchable_count = sum(1 for entry in entries if entry["wordCount"])
        section = {
            "id": item["id"],
            "title": item["title"],
            "kind": item["kind"],
            "material": item["material"],
            "theme": item["theme"],
            "contentLabel": content_label,
            "itemCount": len(entries),
            "searchableCount": searchable_count,
            "entries": entries,
            "summary": sample_text(entries) or item["theme"],
        }
        sections.append(section)

    data = {
        "course": {
            "title": "Simulation and Modeling",
            "subtitle": "Discrete-Event System Simulation",
            "source": "Banks, Carson II, Nelson & Nicol, elaborated by Dr. Farouk Shaaban",
            "lectureMaterial": "materials/Simulation.pdf",
            "pageCount": len(lecture_pages),
            "grading": [
                "4th week lecture test: 10 marks",
                "7th week lecture test: 15 marks",
                "7th week section test: 5 marks",
                "12th week lecture test: 20 marks",
                "5 lecture assignments: 5 marks",
                "Section term work: 5 marks",
                "Final exam: 40 marks",
            ],
        },
        "modules": MODULES,
        "sections": sections,
    }

    output = ROOT / "assets" / "course-data.js"
    output.write_text(
        "window.COURSE_DATA = " + json.dumps(data, ensure_ascii=False, indent=2) + ";\n",
        encoding="utf-8",
    )
    print(f"Wrote {output}")


if __name__ == "__main__":
    main()
